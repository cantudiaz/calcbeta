import matplotlib.pyplot as plt
import pandas as pd
import pandas_datareader.data as web
import datetime
import numpy as np
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches

plt.style.use('ggplot')

# Define the period and tickers of interest.

end = datetime.datetime.today() - relativedelta(months=1)  # last whole month
start = end - relativedelta(months=35)  # for a 36 month Beta
start2 = end - relativedelta(months=11)  # for a 12 month Beta
tickers = ["AMZN", "AAPL", "GOOG", "GOOGL", "IBM", "MSFT", "BRK-B", "FIBRAPL14.MX", "FIBRAMQ12.MX",
           "^MXX", "^GSPC"]


def get_stock_data():  # Gets stock data and resamples it to last data in month.
    df = web.DataReader(tickers, 'yahoo', start, end)
    df = df['Adj Close'].resample('M').last()
    df.to_pickle("stocks.pickle")


def returns_since_start():  # t0 = 1
    df = pd.read_pickle("stocks.pickle")
    df = df.apply(lambda x: x / x[0])
    return df


def percent_change():
    df = pd.read_pickle("stocks.pickle")
    df = df.apply(lambda x: (x - x.shift(1)) / x.shift(1))
    return df


def get_beta():  # obtains betas, st dev, avg return, R^2.
    # Obtain beta
    df = percent_change()
    covariance_matrix = df.cov()
    variance = df.var()
    beta = covariance_matrix / variance
    # Obtain R^2
    r_squared = df.corr('pearson') ** 2
    # Obtain Yearly Std Dev
    standard_error = (variance * 12) ** (1 / 2)
    # Obtain Yearly Avg Returns
    returns_yr = df.mean() * 12
    # Output to a CSV file
    output = pd.DataFrame({"Avg. Return": returns_yr, "Beta IPC": beta["^MXX"],
                           "Std. Error": standard_error, "R^2 IPC": r_squared["^MXX"], "Beta S&P500": beta["^GSPC"],
                           "R^2 S&P500": r_squared["^GSPC"]})
    output = output[["Avg. Return", "Std. Error", "Beta IPC", "R^2 IPC", "Beta S&P500", "R^2 S&P500"]]
    return output


def stocks_plot():  # plots each stock vis a vis the Mexican Index
    df1 = returns_since_start()
    df2 = percent_change()

    for column in df1:
        if column != '^MXX':
            fig = plt.figure()
            ax1 = plt.subplot2grid((2, 1), (0, 0), rowspan=1, colspan=1)
            plt.title("Retornos acumulados y mensuales, " + "{:%m-%Y}".format(start) +
                      " a " "{:%m-%Y}".format(end) + "." )

            ax2 = plt.subplot2grid((2, 1), (1, 0), rowspan=1, colspan=1)

            df1[[column, '^MXX']].plot(ax=ax1).axhline(y=1,
                                                       color="black", lw=2)
            ax1.legend(loc=2)

            df2[[column, '^MXX']].plot(ax=ax2).axhline(y=0,
                                                       color="black", lw=2)
            ax2.legend(loc=2)

            plt.savefig(column + ".png")


def powerpoint_presentation(): #creates a power point presentation and adds slides with graphs
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Gráfica del desenvolvimiento del precio de algunas acciones"
    subtitle.text = "Name"

    prs.save("test.pptx")

    for ticker in sorted(tickers):
        # Add the png files to a slide in a pptx
        if ticker != "^MXX":
            prs = Presentation("test.pptx")
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            left = Inches(1)
            top = Inches(0.5)
            pic = slide.shapes.add_picture(ticker + ".png", left, top)

            prs.save("test.pptx")

# stocks_data = get_stock_data()

stock_returns = returns_since_start()
print(stock_returns)

beta = get_beta()
beta.to_csv("betas.csv")

stocks_plot()

powerpoint_presentation()